"""Анализ файлов: извлечение текста, хеши, определение типов."""

from __future__ import annotations

import hashlib
import mimetypes
import os
import subprocess
from pathlib import Path

from config import ARCHIVE_EXTS, EXECUTABLE_EXTS, IMAGE_EXTS, TEMP_FILE_PATTERNS


AUDIO_EXTS = {"ogg", "mp3", "wav", "flac", "aac", "wma", "m4a", "opus", "aiff"}

# Форматы, которые модель НЕ принимает — надо конвертировать в JPEG
NON_NATIVE_IMAGE_EXTS = {"webp", "bmp", "tiff", "tif", "heic", "heif", "raw", "cr2", "nef", "arw", "svg"}


def is_temp_file(filepath: str) -> bool:
    """Проверить, является ли файл временным/мусорным."""
    import fnmatch
    from pathlib import Path
    name = Path(filepath).name
    for pattern in TEMP_FILE_PATTERNS:
        if fnmatch.fnmatch(name, pattern):
            return True
    return False


def compute_file_hash(filepath: str, max_size: int = 100 * 1024 * 1024) -> str:
    """SHA-256 хеш. Для больших файлов — partial hash."""
    h = hashlib.sha256()
    try:
        size = os.path.getsize(filepath)
        with open(filepath, "rb") as f:
            if size <= max_size:
                while chunk := f.read(8192):
                    h.update(chunk)
            else:
                # Первые и последние 64KB
                h.update(f.read(65536))
                f.seek(size - 65536)
                h.update(f.read(65536))
        return h.hexdigest()
    except Exception as e:
        print(f"[hash] Ошибка {filepath}: {e}")
        return ""


def extract_text(filepath: str) -> str:
    """Извлечь текст из файла."""
    ext = Path(filepath).suffix.lower().lstrip(".")
    size = os.path.getsize(filepath)

    # Маленькие файлы > 10MB не читаем
    if size > 10 * 1024 * 1024:
        return _quick_signature(filepath, ext)

    try:
        if ext in ("txt", "md", "csv", "json", "yaml", "yml", "xml", "log",
                    "cfg", "ini", "conf", "py", "js", "ts", "java", "c", "cpp",
                    "h", "cs", "go", "rs", "sh", "bat", "ps1", "sql", "gpx",
                    "opml", "vcf", "ics", "toml"):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:5000]

        # Office-файлы через python-docx / openpyxl / pptx
        if ext == "docx":
            return _extract_docx(filepath)
        if ext == "xlsx":
            return _extract_xlsx(filepath)
        if ext == "pptx":
            return _extract_pptx(filepath)

        if ext == "pdf":
            return _extract_pdf(filepath)

        if ext == "djvu":
            return _extract_djvu(filepath)

        if ext == "doc":
            return _extract_doc(filepath)

        return _quick_signature(filepath, ext)
    except Exception:
        return _quick_signature(filepath, ext)


def _extract_docx(filepath: str) -> str:
    """Текст из .docx."""
    try:
        import docx
        doc = docx.Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text)[:5000]
    except ImportError:
        # Fallback: zip-распаковка
        return _extract_office_text(filepath)
    except Exception:
        return _extract_office_text(filepath)


def _extract_xlsx(filepath: str) -> str:
    """Текст из .xlsx."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        texts.append(str(cell))
        return "\n".join(texts)[:5000]
    except ImportError:
        return _extract_office_text(filepath)
    except Exception:
        return _extract_office_text(filepath)


def _extract_pptx(filepath: str) -> str:
    """Текст из .pptx включая таблицы."""
    texts = []

    # Пробуем python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_texts.append(cell.text.strip())
                        if row_texts:
                            texts.append(" | ".join(row_texts))
    except ImportError:
        pass
    except Exception:
        pass

    if texts:
        return "\n".join(texts)[:10000]

    # Fallback: zip-распаковка
    return _extract_office_text(filepath)


def _extract_office_text(filepath: str) -> str:
    """Fallback: извлечь текст из Office XML внутри zip."""
    import zipfile
    import re
    try:
        with zipfile.ZipFile(filepath) as zf:
            text = ""
            for name in zf.namelist():
                if not name.endswith(".xml"):
                    continue
                # pptx: slideN.xml, presentation.xml
                # docx: document.xml
                # xlsx: sheetN.xml
                raw = zf.read(name).decode("utf-8", errors="ignore")
                # Вытаскиваем текст из XML-тегов
                clean = re.sub(r"<[^>]+>", " ", raw)
                clean = re.sub(r"\s+", " ", clean).strip()
                if clean:
                    text += clean + "\n"
            return text[:10000] if text else f"[Office файл, размер: {os.path.getsize(filepath)} байт]"
    except Exception:
        return f"[Office файл, размер: {os.path.getsize(filepath)} байт]"


def _extract_pdf(filepath: str) -> str:
    """Текст из PDF."""
    # pdftotext
    try:
        result = subprocess.run(
            ["pdftotext", "-l", "5", "-q", filepath, "-"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            return result.stdout[:5000]
    except (FileNotFoundError, Exception):
        pass

    # pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            text = ""
            for page in pdf.pages[:5]:
                t = page.extract_text()
                if t:
                    text += t + "\n"
            return text[:5000] if text.strip() else ""
    except (ImportError, Exception):
        pass

    return ""


def _extract_djvu(filepath: str) -> str:
    """Извлечь текст из DJVU через djvutxt или ddjvu."""
    try:
        result = subprocess.run(
            ["djvutxt", filepath],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout[:5000]
    except (FileNotFoundError, Exception):
        pass
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_pdf = tmp.name
        subprocess.run(
            ["ddjvu", "-format=pdf", filepath, tmp_pdf],
            capture_output=True, timeout=120,
        )
        if os.path.exists(tmp_pdf):
            text = _extract_pdf(tmp_pdf)
            os.unlink(tmp_pdf)
            return text
    except Exception:
        pass
    return ""


def _extract_doc(filepath: str) -> str:
    """Извлечь текст из .doc (старый формат)."""
    try:
        result = subprocess.run(
            ["antiword", "-m", "UTF-8", filepath],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout[:5000]
    except (FileNotFoundError, Exception):
        pass
    try:
        result = subprocess.run(
            ["catdoc", filepath],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout[:5000]
    except (FileNotFoundError, Exception):
        pass
    return f"[DOC файл, размер: {os.path.getsize(filepath)} байт]"


def pdf_to_images(filepath: str, max_pages: int = 5) -> list[str]:
    """
    Конвертировать PDF-страницы в JPEG-изображения.
    Возвращает список путей к временным файлам.
    Вызывающий обязан удалить их после использования.

    Qwen3.5 мультимодальная модель принимает JPEG/PNG.
    JPEG предпочтительнее — меньше размер при кодировании в base64.
    """
    import tempfile
    images = []
    try:
        # pdftocairo → JPEG (меньше размер, модель принимает)
        with tempfile.TemporaryDirectory(prefix="pdf_") as tmpdir:
            prefix = os.path.join(tmpdir, "page")
            result = subprocess.run(
                [
                    "pdftocairo", "-jpeg", "-r", "150",
                    "-f", "1", "-l", str(max_pages),
                    filepath, prefix
                ],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0:
                for fn in sorted(os.listdir(tmpdir)):
                    if fn.endswith(".jpg"):
                        images.append(os.path.join(tmpdir, fn))
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        # Fallback: pdftoppm
        try:
            with tempfile.TemporaryDirectory(prefix="pdf_") as tmpdir:
                prefix = os.path.join(tmpdir, "page")
                result = subprocess.run(
                    [
                        "pdftoppm", "-jpeg", "-r", "150",
                        "-f", "1", "-l", str(max_pages),
                        filepath, prefix
                    ],
                    capture_output=True, text=True, timeout=120,
                )
                if result.returncode == 0:
                    for fn in sorted(os.listdir(tmpdir)):
                        if fn.endswith(".jpg"):
                            images.append(os.path.join(tmpdir, fn))
        except Exception:
            pass

    return images


def image_to_jpeg(filepath: str, quality: int = 85) -> str:
    """
    Конвертировать изображение в JPEG для совместимости с Qwen3.5.
    Возвращает путь к временному файлу (вызывающий должен удалить).
    Если файл уже JPEG — возвращает исходный путь.
    """
    ext = Path(filepath).suffix.lower().lstrip(".")
    if ext in ("jpg", "jpeg"):
        return filepath  # Уже JPEG

    try:
        from PIL import Image
        tmp_path = filepath + ".converted.jpg"
        img = Image.open(filepath)
        # HEIC/HEIF/RGBA — конвертируем в RGB для JPEG
        if img.mode in ("RGBA", "P", "LA", "PA", "RGBa", "La"):
            # Создаём белый фон для прозрачных изображений
            background = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            if "A" in img.mode:
                background.paste(img, mask=img.split()[-1])
            else:
                background.paste(img)
            img = background
        elif img.mode != "RGB":
            img = img.convert("RGB")
        img.save(tmp_path, "JPEG", quality=quality)
        return tmp_path
    except ImportError:
        # Fallback: ffmpeg через ffprobe+swscale
        try:
            tmp_path = filepath + ".converted.jpg"
            result = subprocess.run(
                ["ffmpeg", "-y", "-i", filepath, "-q:v", "2", tmp_path],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and os.path.exists(tmp_path):
                return tmp_path
        except Exception:
            pass
    except Exception:
        pass

    return filepath  # Вернём как есть, если конвертация не удалась


def _quick_signature(filepath: str, ext: str) -> str:
    """Быстрое описание для файлов без содержимого."""
    size = os.path.getsize(filepath)
    size_str = _human_size(size)
    return f"[{ext.upper()} файл, {size_str}]"


def _human_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} TB"


def is_archive(filepath: str) -> bool:
    return Path(filepath).suffix.lower().lstrip(".") in ARCHIVE_EXTS


def is_executable(filepath: str) -> bool:
    return Path(filepath).suffix.lower().lstrip(".") in EXECUTABLE_EXTS


def is_image(filepath: str) -> bool:
    ext = Path(filepath).suffix.lower().lstrip(".")
    if ext in IMAGE_EXTS:
        return True
    mime, _ = mimetypes.guess_type(filepath)
    return mime is not None and mime.startswith("image/")
